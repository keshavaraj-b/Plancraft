"""
PlanCraft — Code Games 2026 Presentation Generator

Creates a professional PowerPoint presentation for the Code Games 2026 competition.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

# Wolters Kluwer brand colors
WK_DARK_BLUE = RGBColor(0x00, 0x2B, 0x49)  # Primary dark
WK_TEAL = RGBColor(0x00, 0x7B, 0x8A)       # Accent teal
WK_LIGHT_BLUE = RGBColor(0x00, 0xA3, 0xC4)  # Light accent
WK_GREEN = RGBColor(0x6C, 0xBE, 0x45)       # Success/highlight
WK_ORANGE = RGBColor(0xF2, 0x8C, 0x28)      # Attention
WK_GRAY = RGBColor(0x5A, 0x5A, 0x5A)        # Body text
WK_LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)  # Background
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)


def set_slide_bg(slide, color):
    """Set solid background color for a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_styled_textbox(slide, left, top, width, height, text, font_size=18,
                       bold=False, color=WK_GRAY, alignment=PP_ALIGN.LEFT,
                       font_name="Segoe UI"):
    """Add a styled text box to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_points(slide, left, top, width, height, items, font_size=16,
                      color=WK_GRAY, bullet_color=WK_TEAL):
    """Add bullet point list to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Segoe UI"
        p.space_after = Pt(8)
        p.level = 0
    return txBox


def add_section_header(slide, title, subtitle=None):
    """Add a consistent section header to a slide."""
    set_slide_bg(slide, WHITE)

    # Top accent bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.08))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WK_TEAL
    shape.line.fill.background()

    # Title
    add_styled_textbox(slide, Inches(0.8), Inches(0.5), Inches(8.4), Inches(1),
                       title, font_size=32, bold=True, color=WK_DARK_BLUE)

    if subtitle:
        add_styled_textbox(slide, Inches(0.8), Inches(1.3), Inches(8.4), Inches(0.6),
                           subtitle, font_size=16, color=WK_GRAY)

    # Bottom separator line
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.9), Inches(2), Inches(0.04))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WK_TEAL
    shape.line.fill.background()


def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)  # 16:9 aspect ratio

    # =========================================================================
    # SLIDE 1: Title Slide
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    set_slide_bg(slide, WK_DARK_BLUE)

    # Decorative accent
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.12))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WK_TEAL
    shape.line.fill.background()

    # Title
    add_styled_textbox(slide, Inches(0.8), Inches(1.2), Inches(8.4), Inches(1.2),
                       "PlanCraft", font_size=48, bold=True, color=WHITE)

    # Subtitle
    add_styled_textbox(slide, Inches(0.8), Inches(2.2), Inches(8.4), Inches(0.8),
                       "AI Planning Agent — SOW to Jira-Ready Stories in Minutes",
                       font_size=20, color=WK_LIGHT_BLUE)

    # Event
    add_styled_textbox(slide, Inches(0.8), Inches(3.5), Inches(8.4), Inches(0.5),
                       "Code Games 2026 | Wolters Kluwer",
                       font_size=14, color=WK_LIGHT_BLUE)

    # Bottom accent
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(5.5), Inches(10), Inches(0.125))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WK_GREEN
    shape.line.fill.background()

    # =========================================================================
    # SLIDE 2: Team Introduction
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "Team Introduction")

    add_bullet_points(slide, Inches(0.8), Inches(2.2), Inches(8.4), Inches(3),
                      [
                          "Team: [Your Team Name]",
                          "Members: [Member 1], [Member 2], [Member 3], [Member 4]",
                          "Division: Legal Software — Passport Product Line",
                          "Roles: Full-Stack Dev, Platform Engineer, QA, Product",
                      ], font_size=18)

    # =========================================================================
    # SLIDE 3: Advisor Acknowledgment
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "Advisor Acknowledgment")

    add_bullet_points(slide, Inches(0.8), Inches(2.2), Inches(8.4), Inches(3),
                      [
                          "Advisor: [Advisor Name]",
                          "Role: [Advisor's Title / Department]",
                          "Guidance provided on: Architecture design, Passport domain expertise,",
                          "    and MCP protocol integration patterns",
                          "",
                          "Thank you for the mentorship and technical direction!",
                      ], font_size=18)

    # =========================================================================
    # SLIDE 4: Project Title & Summary
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "Project: PlanCraft",
                       "AI Planning Agent for Statement of Work → Development Stories")

    add_bullet_points(slide, Inches(0.8), Inches(2.3), Inches(8.4), Inches(3),
                      [
                          "Problem: Manual SOW analysis takes days/weeks, misses existing features",
                          "Solution: AI agent that parses SOWs, queries live Passport knowledge,",
                          "    and generates Jira-ready epics/stories with zero LLM cost",
                          "How: MCP server integrated into VS Code Copilot Chat — uses",
                          "    Passport Copilot's 113 tools for real domain intelligence",
                      ], font_size=16)

    # =========================================================================
    # SLIDE 5: The Problem
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "The Problem", "Why SOW → Stories is painful today")

    add_bullet_points(slide, Inches(0.8), Inches(2.3), Inches(4.2), Inches(3),
                      [
                          "Manual analysis of 50+ page SOWs",
                          "No awareness of existing Passport features",
                          "Duplicate work — building what already exists",
                          "Inconsistent story quality across teams",
                          "No client-specific context (customizations,",
                          "    business rules, Groovy scripts)",
                      ], font_size=14)

    # Right side - impact numbers
    add_styled_textbox(slide, Inches(5.5), Inches(2.3), Inches(4), Inches(0.5),
                       "Impact:", font_size=16, bold=True, color=WK_DARK_BLUE)
    add_bullet_points(slide, Inches(5.5), Inches(2.8), Inches(4), Inches(2.5),
                      [
                          "2-3 weeks → minutes for story generation",
                          "30%+ stories duplicate existing features",
                          "Inconsistent acceptance criteria",
                          "Missing data migration tasks",
                      ], font_size=14, color=WK_ORANGE)

    # =========================================================================
    # SLIDE 6: Architecture Overview
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "Architecture", "How PlanCraft works end-to-end")

    # Architecture diagram as text (clean representation)
    arch_text = (
        "┌─────────────────────────────────────────────────────────┐\n"
        "│          VS Code + GitHub Copilot (LLM)                 │\n"
        "│                      │                                  │\n"
        "│              PlanCraft MCP Server                        │\n"
        "│    ┌─────────┬──────────┬────────────┬─────────┐        │\n"
        "│    │parse_sow│get_context│detect_overlap│validate│        │\n"
        "│    └─────────┴──────────┴────────────┴─────────┘        │\n"
        "│                      │                                  │\n"
        "│         Passport Copilot MCP (113 tools)                │\n"
        "│    ┌──────────┬──────────┬───────────┬──────────┐       │\n"
        "│    │ Schema   │ Docs     │ Workflows │ Rules    │       │\n"
        "│    └──────────┴──────────┴───────────┴──────────┘       │\n"
        "└─────────────────────────────────────────────────────────┘"
    )
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(8.4), Inches(3.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = arch_text
    p.font.size = Pt(10)
    p.font.name = "Consolas"
    p.font.color.rgb = WK_DARK_BLUE

    # =========================================================================
    # SLIDE 7: Key Innovation — Zero LLM Cost
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "Key Innovation", "GitHub Copilot IS the LLM — Zero API Cost")

    add_bullet_points(slide, Inches(0.8), Inches(2.3), Inches(4.2), Inches(3),
                      [
                          "No OpenAI/Azure API keys needed",
                          "No token cost per generation",
                          "Copilot handles all LLM reasoning",
                          "PlanCraft provides tools + context",
                          "MCP protocol = standard interface",
                      ], font_size=15)

    # Right column - comparison
    add_styled_textbox(slide, Inches(5.5), Inches(2.3), Inches(4), Inches(0.5),
                       "vs. Traditional Approach:", font_size=14, bold=True, color=WK_DARK_BLUE)
    add_bullet_points(slide, Inches(5.5), Inches(2.8), Inches(4), Inches(2.5),
                      [
                          "❌ Custom LLM pipeline = $$$/month",
                          "❌ Prompt engineering maintenance",
                          "❌ Token limits on large SOWs",
                          "✅ PlanCraft = free with Copilot license",
                          "✅ Copilot improves → PlanCraft improves",
                      ], font_size=13)

    # =========================================================================
    # SLIDE 8: Live Passport Integration
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "Live Passport Intelligence",
                       "Real data, real schema, real client customizations")

    add_bullet_points(slide, Inches(0.8), Inches(2.3), Inches(8.4), Inches(3.2),
                      [
                          "Spawns Passport Copilot as MCP subprocess (stdio protocol)",
                          "Connects to specific client DB (e.g., CITI-D3, AMEX-P1)",
                          "Queries 113 real tools: schema, docs, workflows, Groovy scripts",
                          "Discovers existing customizations before generating stories",
                          "Runs impact analysis to flag high-risk changes",
                          "Result: Stories that reference actual Passport entities, tables, and rules",
                      ], font_size=15)

    # =========================================================================
    # SLIDE 9: Workflow Demo Introduction
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "Demo: End-to-End Workflow",
                       "From SOW document to validated Jira stories in Copilot Chat")

    # Steps
    steps = [
        "1️⃣  Select PlanCraft agent in VS Code Copilot Chat",
        "2️⃣  Paste SOW text or attach PDF/DOCX → parse_sow extracts concepts",
        "3️⃣  get_passport_context connects to client, gathers domain knowledge",
        "4️⃣  detect_feature_overlap finds what already exists in Passport",
        "5️⃣  Copilot generates structured epics/features/stories using context",
        "6️⃣  validate_stories checks against Passport rules + impact analysis",
        "7️⃣  Export to Word, PDF, or Jira-preview format",
    ]
    add_bullet_points(slide, Inches(0.8), Inches(2.2), Inches(8.4), Inches(3.2),
                      steps, font_size=14)

    # =========================================================================
    # SLIDE 10: Demo - SOW Parsing
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "Demo: SOW Parsing", "Intelligent concept extraction from documents")

    # Show what parse_sow does
    add_bullet_points(slide, Inches(0.8), Inches(2.3), Inches(4.5), Inches(3),
                      [
                          "Supports PDF, DOCX, and plain text",
                          "Extracts key concepts via regex patterns",
                          "Identifies Passport entities mentioned",
                          "Detects modules referenced in SOW",
                          "Caps at 30 concepts to avoid overload",
                      ], font_size=14)

    # Output example
    add_styled_textbox(slide, Inches(5.5), Inches(2.3), Inches(4.2), Inches(0.4),
                       "Output:", font_size=12, bold=True, color=WK_DARK_BLUE)
    output_text = (
        '{\n'
        '  "concepts": ["Invoice Approval", ...],\n'
        '  "entities": ["matter", "invoice", ...],\n'
        '  "modules": ["eBilling", "Workflow"]\n'
        '}'
    )
    txBox = slide.shapes.add_textbox(Inches(5.5), Inches(2.7), Inches(4.2), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = output_text
    p.font.size = Pt(11)
    p.font.name = "Consolas"
    p.font.color.rgb = WK_DARK_BLUE

    # =========================================================================
    # SLIDE 11: Demo - Story Output
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "Demo: Generated Story Output", "Jira-ready structured JSON")

    story_example = (
        '{\n'
        '  "title": "As a legal ops manager, I want automated\n'
        '           invoice routing, so that approvals happen\n'
        '           within SLA",\n'
        '  "acceptance_criteria": [\n'
        '    "Given an invoice > $10K, When submitted,\n'
        '     Then route to senior approver"\n'
        '  ],\n'
        '  "tasks": ["Configure workflow rule in Groovy",\n'
        '            "Add approval threshold to business rules"],\n'
        '  "story_points": 5,\n'
        '  "priority": "High"\n'
        '}'
    )
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(8.4), Inches(3.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = story_example
    p.font.size = Pt(12)
    p.font.name = "Consolas"
    p.font.color.rgb = WK_DARK_BLUE

    # =========================================================================
    # SLIDE 12: Feature Overlap Detection
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "Smart Feature Detection",
                       "Avoids duplicating what Passport already provides")

    add_bullet_points(slide, Inches(0.8), Inches(2.3), Inches(8.4), Inches(3.2),
                      [
                          "Searches docs, code, workflows, and screens for each concept",
                          "Recommendations: skip_or_extend | extend_or_customize | new_development",
                          "Prevents teams from building what's already configurable",
                          "Example: 'Invoice approval routing' → Already exists as workflow!",
                          "    → Recommendation: 'extend_or_customize' with Groovy config task",
                          "",
                          "Result: 30%+ reduction in wasted development effort",
                      ], font_size=14)

    # =========================================================================
    # SLIDE 13: Value Explanation
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "Business Value", "Why PlanCraft matters")

    # Left column
    add_styled_textbox(slide, Inches(0.8), Inches(2.2), Inches(4), Inches(0.4),
                       "Time Savings", font_size=16, bold=True, color=WK_TEAL)
    add_bullet_points(slide, Inches(0.8), Inches(2.6), Inches(4), Inches(1.2),
                      [
                          "2-3 weeks → 30 minutes",
                          "Instant domain knowledge lookup",
                          "Automated validation + impact analysis",
                      ], font_size=13)

    add_styled_textbox(slide, Inches(0.8), Inches(3.8), Inches(4), Inches(0.4),
                       "Quality", font_size=16, bold=True, color=WK_TEAL)
    add_bullet_points(slide, Inches(0.8), Inches(4.2), Inches(4), Inches(1.2),
                      [
                          "Consistent story format (As a / Given-When-Then)",
                          "No missed migration or config tasks",
                          "Client-specific, not generic",
                      ], font_size=13)

    # Right column
    add_styled_textbox(slide, Inches(5.5), Inches(2.2), Inches(4), Inches(0.4),
                       "Cost", font_size=16, bold=True, color=WK_TEAL)
    add_bullet_points(slide, Inches(5.5), Inches(2.6), Inches(4), Inches(1.2),
                      [
                          "Zero LLM cost (uses existing Copilot license)",
                          "Reduces rework from duplicate features",
                          "Fewer defects from better acceptance criteria",
                      ], font_size=13)

    add_styled_textbox(slide, Inches(5.5), Inches(3.8), Inches(4), Inches(0.4),
                       "Scalability", font_size=16, bold=True, color=WK_TEAL)
    add_bullet_points(slide, Inches(5.5), Inches(4.2), Inches(4), Inches(1.2),
                      [
                          "Open plugin architecture for any WK product",
                          "T360, HighQ, etc. — just add MCP server",
                          "One tool, all product lines",
                      ], font_size=13)

    # =========================================================================
    # SLIDE 14: What Was Built During Code Games
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "Built During Code Games 2026",
                       "Everything was created from scratch in this competition")

    add_bullet_points(slide, Inches(0.8), Inches(2.3), Inches(4.5), Inches(3.2),
                      [
                          "✅ FastMCP server (Python, MCP protocol)",
                          "✅ SOW parser (PDF, DOCX, text)",
                          "✅ Concept extraction engine",
                          "✅ Passport Copilot MCP client integration",
                          "✅ Feature overlap detection",
                          "✅ Story validation with impact analysis",
                          "✅ Export: Word, PDF, Jira-preview",
                          "✅ VS Code agent definition + config",
                      ], font_size=14)

    # Right side - tech stack
    add_styled_textbox(slide, Inches(5.8), Inches(2.3), Inches(3.8), Inches(0.4),
                       "Tech Stack:", font_size=14, bold=True, color=WK_DARK_BLUE)
    add_bullet_points(slide, Inches(5.8), Inches(2.7), Inches(3.8), Inches(2.8),
                      [
                          "Python 3.11+ / FastMCP",
                          "MCP Protocol (stdio transport)",
                          "VS Code + GitHub Copilot",
                          "PyPDF2, python-docx, fpdf2",
                          "Passport Copilot (113 tools)",
                          "Pydantic for validation",
                      ], font_size=13)

    # =========================================================================
    # SLIDE 15: AI Usage Disclosure
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "AI Usage Disclosure")

    add_bullet_points(slide, Inches(0.8), Inches(2.3), Inches(8.4), Inches(3),
                      [
                          "GitHub Copilot was used as a coding assistant during development",
                          "GitHub Copilot IS the runtime LLM — PlanCraft is built to leverage it",
                          "No external LLM APIs (OpenAI, Azure, etc.) are called at runtime",
                          "All architecture decisions, integration patterns, and tool design",
                          "    were made by the team — AI assisted with implementation speed",
                          "",
                          "AI is both our development tool AND our product's engine",
                      ], font_size=16)

    # =========================================================================
    # SLIDE 16: Future Roadmap
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "Future Roadmap", "Where PlanCraft goes next")

    add_bullet_points(slide, Inches(0.8), Inches(2.3), Inches(4.5), Inches(3),
                      [
                          "Phase 2:",
                          "  • Direct Jira API push (not just preview)",
                          "  • Sprint capacity planning",
                          "  • Dependency graph between stories",
                          "  • Test case auto-generation",
                      ], font_size=14)

    add_bullet_points(slide, Inches(5.5), Inches(2.3), Inches(4.2), Inches(3),
                      [
                          "Phase 3:",
                          "  • T360 agent integration",
                          "  • HighQ agent integration",
                          "  • Cross-product story planning",
                          "  • SOW comparison / change tracking",
                      ], font_size=14)

    # =========================================================================
    # SLIDE 17: Thank You / Q&A
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WK_DARK_BLUE)

    # Top accent
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.12))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WK_TEAL
    shape.line.fill.background()

    add_styled_textbox(slide, Inches(0.8), Inches(1.5), Inches(8.4), Inches(1),
                       "Thank You!", font_size=44, bold=True, color=WHITE,
                       alignment=PP_ALIGN.CENTER)

    add_styled_textbox(slide, Inches(0.8), Inches(2.8), Inches(8.4), Inches(0.6),
                       "PlanCraft — AI-Powered SOW to Stories",
                       font_size=20, color=WK_LIGHT_BLUE, alignment=PP_ALIGN.CENTER)

    add_styled_textbox(slide, Inches(0.8), Inches(3.8), Inches(8.4), Inches(0.5),
                       "Questions?",
                       font_size=18, color=WK_GREEN, alignment=PP_ALIGN.CENTER)

    add_styled_textbox(slide, Inches(0.8), Inches(4.6), Inches(8.4), Inches(0.5),
                       "Code Games 2026 | [Team Name]",
                       font_size=12, color=WK_LIGHT_BLUE, alignment=PP_ALIGN.CENTER)

    # Bottom accent
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(5.5), Inches(10), Inches(0.125))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WK_GREEN
    shape.line.fill.background()

    # =========================================================================
    # Save
    # =========================================================================
    output_dir = Path("output")
    output_dir.mkdir(parents=True, exist_ok=True)
    output_path = output_dir / "PlanCraft_CodeGames2026_Presentation.pptx"
    prs.save(str(output_path))
    print(f"✅ Presentation saved to: {output_path}")
    print(f"   Total slides: {len(prs.slides)}")
    return str(output_path)


if __name__ == "__main__":
    create_presentation()
